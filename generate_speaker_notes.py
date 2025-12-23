import io
import os
import subprocess
import sys
import base64
import tempfile
from typing import List, Tuple
import fitz  # PyMuPDF
from pptx import Presentation
from anthropic import Anthropic


class PPTXReader:
    def pptx_to_pdf(self, file_path: str, output_path: str) -> str:
        """LibreOfficeを使用してPPTXをPDFに変換"""
        subprocess.run([
            'libreoffice', '--headless', '--convert-to', 'pdf', 
            '--outdir', os.path.dirname(output_path), file_path
        ], check=True)
        return output_path


class SpeakerNotesGenerator:
    def __init__(self, api_key: str = None):
        """
        発表者ノート生成器
        
        Args:
            api_key: Anthropic APIキー (Noneの場合は環境変数から取得)
        """
        if api_key is None:
            api_key = os.environ.get("ANTHROPIC_API_KEY")
        
        if not api_key:
            raise ValueError("APIキーが設定されていません。環境変数ANTHROPIC_API_KEYを設定するか、引数で渡してください。")
        
        self.api_key = api_key
        self.pptx_reader = PPTXReader()
        self.client = Anthropic(api_key=api_key)

    def pdf_to_images(self, pdf_path: str) -> List[bytes]:
        """
        PDFファイルを画像データのリストに変換
        
        Args:
            pdf_path: PDFファイルのパス
            
        Returns:
            List[bytes]: 各ページの画像データ（PNG形式）
        """
        doc = fitz.open(pdf_path)
        images = []
        
        try:
            for page_num in range(len(doc)):
                page = doc.load_page(page_num)
                # 高解像度で画像を生成（2倍ズーム）
                pix = page.get_pixmap(matrix=fitz.Matrix(2.0, 2.0))
                img_data = pix.tobytes("png")
                images.append(img_data)
        finally:
            doc.close()
        
        return images

    def generate_script_from_image(self, image_data: bytes, slide_number: int) -> str:
        """
        スライド画像から発表者原稿を生成
        
        Args:
            image_data: スライドの画像データ（PNG形式）
            slide_number: スライド番号（表示用）
            
        Returns:
            str: 生成された発表者原稿
        """
        # 画像をbase64エンコード
        image_base64 = base64.b64encode(image_data).decode()
        
        prompt = f"""このスライド（{slide_number}ページ目）の内容を元に、発表者向けの詳しい原稿を日本語で作成してください。

要件:
- 聞き手にとって分かりやすい説明になるよう、具体例や補足説明を含める
- スライドに書かれていない背景情報や詳細説明も追加する
- 自然な話し言葉で、実際の発表で使いやすい形式にする
- 2-3分程度で話せる分量（400-600文字程度）にする
- スライドの内容を単純に読み上げるのではなく、聞き手に分かりやすく説明する

形式:
原稿のみを出力してください。「このスライドでは」などの前置きは不要です。"""

        try:
            message = self.client.messages.create(
                model="claude-sonnet-4-5-20250929",
                max_tokens=4000,
                messages=[{
                    "role": "user",
                    "content": [
                        {
                            "type": "image", 
                            "source": {
                                "type": "base64", 
                                "media_type": "image/png", 
                                "data": image_base64
                            }
                        },
                        {"type": "text", "text": prompt}
                    ]
                }]
            )
            return message.content[0].text.strip()
        except Exception as e:
            print(f"スライド {slide_number} の原稿生成でエラー: {e}")
            return f"スライド {slide_number} の原稿を生成できませんでした。"

    def add_notes_to_pptx(self, pptx_path: str, notes_list: List[str], output_path: str):
        """
        PPTXファイルに発表者ノートを追加
        
        Args:
            pptx_path: 元のPPTXファイルのパス
            notes_list: 各スライドの発表者ノート
            output_path: 出力PPTXファイルのパス
        """
        prs = Presentation(pptx_path)
        
        for i, note_text in enumerate(notes_list):
            if i < len(prs.slides):
                slide = prs.slides[i]
                notes_slide = slide.notes_slide
                text_frame = notes_slide.notes_text_frame
                
                # 既存のノートがある場合は追記、ない場合は新規作成
                if text_frame.text.strip():
                    text_frame.text += "\n\n" + note_text
                else:
                    text_frame.text = note_text
        
        prs.save(output_path)

    def process_pptx_with_pdf(self, input_pptx_path: str, input_pdf_path: str,
                              output_pptx_path: str, progress_callback=None) -> Tuple[bool, str, List[str]]:
        """
        PPTXファイルとPDFファイルから発表者ノートを生成
        （PPTX→PDF変換をスキップ）

        Args:
            input_pptx_path: 入力PPTXファイルのパス
            input_pdf_path: 入力PDFファイルのパス（スライドの画像抽出用）
            output_pptx_path: 出力PPTXファイルのパス
            progress_callback: 進捗コールバック関数 (step_name, current, total)

        Returns:
            Tuple[bool, str, List[str]]: (成功フラグ, メッセージ, 生成された原稿リスト)
        """
        notes_list = []

        try:
            # 1. PDF → 画像変換
            if progress_callback:
                progress_callback("PDF → 画像変換中...", 0, 3)

            images = self.pdf_to_images(input_pdf_path)
            total_slides = len(images)

            # 2. 各スライドの原稿生成
            if progress_callback:
                progress_callback(f"スライド原稿生成中...", 1, 3)

            for i, image_data in enumerate(images):
                print(f"スライド {i+1}/{total_slides} の原稿を生成中...")
                script = self.generate_script_from_image(image_data, i+1)
                notes_list.append(script)

            # 3. PPTXに発表者ノート追加
            if progress_callback:
                progress_callback("PPTX ファイル更新中...", 2, 3)

            self.add_notes_to_pptx(input_pptx_path, notes_list, output_pptx_path)

            if progress_callback:
                progress_callback("完了", 3, 3)

            return True, f"発表者ノート生成完了！ {total_slides}枚のスライドを処理しました。", notes_list

        except Exception as e:
            error_msg = f"処理中にエラーが発生しました: {str(e)}"
            print(error_msg)
            return False, error_msg, notes_list

    def process_pptx(self, input_pptx_path: str, output_pptx_path: str,
                     progress_callback=None) -> Tuple[bool, str, List[str]]:
        """
        PPTXファイル全体を処理して発表者ノートを生成

        Args:
            input_pptx_path: 入力PPTXファイルのパス
            output_pptx_path: 出力PPTXファイルのパス
            progress_callback: 進捗コールバック関数 (step_name, current, total)

        Returns:
            Tuple[bool, str, List[str]]: (成功フラグ, メッセージ, 生成された原稿リスト)
        """
        notes_list = []

        try:
            # 1. PPTX → PDF変換
            if progress_callback:
                progress_callback("PPTX → PDF変換中...", 0, 4)

            with tempfile.TemporaryDirectory() as temp_dir:
                pdf_path = os.path.join(temp_dir, "slides.pdf")
                self.pptx_reader.pptx_to_pdf(input_pptx_path, pdf_path)

                # 2. PDF → 画像変換
                if progress_callback:
                    progress_callback("PDF → 画像変換中...", 1, 4)

                images = self.pdf_to_images(pdf_path)
                total_slides = len(images)

                # 3. 各スライドの原稿生成
                if progress_callback:
                    progress_callback(f"スライド原稿生成中...", 2, 4)

                for i, image_data in enumerate(images):
                    print(f"スライド {i+1}/{total_slides} の原稿を生成中...")
                    script = self.generate_script_from_image(image_data, i+1)
                    notes_list.append(script)

                # 4. PPTXに発表者ノート追加
                if progress_callback:
                    progress_callback("PPTX ファイル更新中...", 3, 4)

                self.add_notes_to_pptx(input_pptx_path, notes_list, output_pptx_path)

                if progress_callback:
                    progress_callback("完了", 4, 4)

                return True, f"発表者ノート生成完了！ {total_slides}枚のスライドを処理しました。", notes_list

        except subprocess.CalledProcessError as e:
            error_msg = f"LibreOfficeでのPDF変換に失敗しました: {e}"
            print(error_msg)
            return False, error_msg, notes_list
        except Exception as e:
            error_msg = f"処理中にエラーが発生しました: {str(e)}"
            print(error_msg)
            return False, error_msg, notes_list


if __name__ == "__main__":
    # テスト実行
    if len(sys.argv) < 3:
        print("使用方法: python generate_speaker_notes.py <input.pptx> <output.pptx>")
        sys.exit(1)
    
    input_file = sys.argv[1]
    output_file = sys.argv[2]
    
    try:
        generator = SpeakerNotesGenerator()
        
        def print_progress(step_name, current, total):
            print(f"[{current}/{total}] {step_name}")
        
        success, message, notes = generator.process_pptx(
            input_file, output_file, progress_callback=print_progress
        )
        
        print(f"\n結果: {message}")
        
        if success:
            print(f"出力ファイル: {output_file}")
            print(f"生成された原稿数: {len(notes)}")
        
    except ValueError as e:
        print(f"エラー: {e}")
        sys.exit(1)