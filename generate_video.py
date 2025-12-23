"""
講義動画生成モジュール
発表者ノート付きPPTXとスライド画像ZIPから、発表者ノートを読み上げる講義動画を生成する
"""

import os
import re
import subprocess
import tempfile
import wave
import zipfile
from pathlib import Path
from typing import List, Optional, Tuple, Callable

from PIL import Image
from pptx import Presentation
from google import genai
from google.genai import types


class VideoGenerator:
    """講義動画生成クラス"""

    def __init__(self, google_api_key: str):
        """
        初期化

        Args:
            google_api_key: Google APIキー（Gemini TTS用）
        """
        if not google_api_key:
            raise ValueError("Google APIキーが設定されていません")

        self.client = genai.Client(api_key=google_api_key)

    def extract_speaker_notes(self, pptx_path: str) -> List[Optional[str]]:
        """
        PPTXファイルから発表者ノートを抽出

        Args:
            pptx_path: PPTXファイルのパス

        Returns:
            List[Optional[str]]: 各スライドの発表者ノート（ない場合はNone）
        """
        presentation = Presentation(pptx_path)
        speaker_notes = []

        for slide in presentation.slides:
            note_text = None

            if slide.has_notes_slide:
                notes_slide = slide.notes_slide
                if notes_slide.notes_text_frame and notes_slide.notes_text_frame.text:
                    note_text = notes_slide.notes_text_frame.text.strip()
                    if not note_text:
                        note_text = None

            speaker_notes.append(note_text)

        return speaker_notes

    def extract_images_from_zip(self, zip_path: str) -> List[Image.Image]:
        """
        ZIPファイルからスライド画像を抽出（スライド番号順）

        ファイル名形式: スライド1.jpeg, スライド2.jpeg, ...

        Args:
            zip_path: ZIPファイルのパス

        Returns:
            List[Image.Image]: スライド番号順の画像リスト
        """
        images = []

        with zipfile.ZipFile(zip_path, 'r') as zf:
            # ファイル名からスライド番号を抽出してソート
            jpeg_files = []
            for name in zf.namelist():
                # ディレクトリはスキップ
                if name.endswith('/'):
                    continue
                # macOSのメタデータフォルダをスキップ
                if '__MACOSX' in name:
                    continue
                # 隠しファイルをスキップ
                basename = os.path.basename(name)
                if basename.startswith('.'):
                    continue
                # JPEG/JPGファイルのみ対象
                if name.lower().endswith(('.jpeg', '.jpg')):
                    # ファイル名から数字を抽出
                    match = re.search(r'(\d+)', basename)
                    if match:
                        num = int(match.group(1))
                        jpeg_files.append((num, name))

            # スライド番号順にソート
            jpeg_files.sort(key=lambda x: x[0])

            # 画像を読み込み
            for _, name in jpeg_files:
                with zf.open(name) as f:
                    image = Image.open(f)
                    image.load()  # メモリに読み込み
                    images.append(image.copy())

        return images

    def text_to_speech(self, text: str, output_path: str) -> str:
        """
        Gemini TTSでテキストを音声に変換

        Args:
            text: 読み上げるテキスト
            output_path: 出力WAVファイルのパス

        Returns:
            str: 出力ファイルのパス
        """
        response = self.client.models.generate_content(
            model="gemini-2.5-flash-preview-tts",
            contents=text,
            config=types.GenerateContentConfig(
                response_modalities=["AUDIO"],
                speech_config=types.SpeechConfig(
                    voice_config=types.VoiceConfig(
                        prebuilt_voice_config=types.PrebuiltVoiceConfig(
                            voice_name="Kore",
                        )
                    )
                ),
            )
        )

        audio_data = response.candidates[0].content.parts[0].inline_data.data

        with wave.open(output_path, "wb") as wf:
            wf.setnchannels(1)
            wf.setsampwidth(2)
            wf.setframerate(24000)
            wf.writeframes(audio_data)

        return output_path

    def create_slide_video(
        self,
        image: Image.Image,
        audio_path: Optional[str],
        output_path: str,
        duration: float = 3.0
    ) -> str:
        """
        スライド画像と音声から動画を作成

        Args:
            image: スライド画像
            audio_path: 音声ファイルのパス（Noneの場合は無音動画）
            output_path: 出力動画ファイルのパス
            duration: 無音動画の場合の長さ（秒）

        Returns:
            str: 出力ファイルのパス
        """
        # 画像を一時ファイルに保存（サイズを偶数に調整）
        with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as temp_file:
            width, height = image.size
            if height % 2 != 0:
                height -= 1
            if width % 2 != 0:
                width -= 1

            resized_image = image.resize((width, height), Image.LANCZOS)
            resized_image.save(temp_file.name)
            temp_image_path = temp_file.name

        try:
            if audio_path:
                # 音声の長さを取得
                duration_cmd = [
                    'ffprobe', '-v', 'error', '-show_entries', 'format=duration',
                    '-of', 'default=noprint_wrappers=1:nokey=1', audio_path
                ]
                duration_result = subprocess.run(duration_cmd, capture_output=True, text=True)
                duration = float(duration_result.stdout.strip())

                # 音声付き動画を作成
                ffmpeg_cmd = [
                    'ffmpeg', '-y',
                    '-loop', '1',
                    '-i', temp_image_path,
                    '-i', audio_path,
                    '-c:v', 'libx264',
                    '-c:a', 'aac',
                    '-t', str(duration),
                    '-pix_fmt', 'yuv420p',
                    '-r', '1',
                    '-shortest',
                    output_path
                ]
            else:
                # 無音動画を作成
                ffmpeg_cmd = [
                    'ffmpeg', '-y',
                    '-loop', '1',
                    '-i', temp_image_path,
                    '-f', 'lavfi',
                    '-i', 'anullsrc=channel_layout=mono:sample_rate=48000',
                    '-c:v', 'libx264',
                    '-c:a', 'aac',
                    '-t', str(duration),
                    '-pix_fmt', 'yuv420p',
                    '-r', '1',
                    '-shortest',
                    output_path
                ]

            result = subprocess.run(ffmpeg_cmd, capture_output=True, text=True)
            if result.returncode != 0:
                raise subprocess.CalledProcessError(result.returncode, ffmpeg_cmd, result.stderr)

            return output_path

        finally:
            os.unlink(temp_image_path)

    def combine_videos(self, video_paths: List[str], output_path: str) -> str:
        """
        複数の動画を結合

        Args:
            video_paths: 結合する動画ファイルのパスリスト
            output_path: 出力ファイルのパス

        Returns:
            str: 出力ファイルのパス
        """
        with tempfile.NamedTemporaryFile(mode='w', suffix='.txt', delete=False) as list_file:
            for video_path in video_paths:
                list_file.write(f"file '{os.path.abspath(video_path)}'\n")
            list_file_path = list_file.name

        try:
            ffmpeg_cmd = [
                'ffmpeg', '-y',
                '-f', 'concat',
                '-safe', '0',
                '-i', list_file_path,
                '-c', 'copy',
                output_path
            ]

            result = subprocess.run(ffmpeg_cmd, capture_output=True, text=True)
            if result.returncode != 0:
                raise subprocess.CalledProcessError(result.returncode, ffmpeg_cmd, result.stderr)

            return output_path

        finally:
            os.unlink(list_file_path)

    def generate_video(
        self,
        pptx_path: str,
        zip_path: str,
        output_path: str,
        progress_callback: Optional[Callable[[str, int, int], None]] = None
    ) -> Tuple[bool, str]:
        """
        講義動画を生成するメイン処理

        Args:
            pptx_path: PPTXファイルのパス（発表者ノート付き）
            zip_path: スライド画像ZIPファイルのパス
            output_path: 出力動画ファイルのパス
            progress_callback: 進捗コールバック関数 (step_name, current, total)

        Returns:
            Tuple[bool, str]: (成功フラグ, メッセージ)
        """
        try:
            # 1. 発表者ノートを抽出
            if progress_callback:
                progress_callback("発表者ノート抽出中...", 0, 1)

            speaker_notes = self.extract_speaker_notes(pptx_path)

            # 2. ZIPから画像を抽出
            if progress_callback:
                progress_callback("スライド画像抽出中...", 0, 1)

            images = self.extract_images_from_zip(zip_path)
            total_slides = len(images)

            if total_slides == 0:
                return False, "ZIPファイルにスライド画像が見つかりません（スライド1.jpeg, スライド2.jpeg...の形式で保存してください）"

            if len(speaker_notes) != total_slides:
                return False, f"スライド数が一致しません（PPTX: {len(speaker_notes)}枚, ZIP: {total_slides}枚）"

            total_steps = total_slides + 1  # スライド数 + 動画結合

            # 一時ディレクトリを作成
            with tempfile.TemporaryDirectory() as temp_dir:
                video_paths = []

                # 3. 各スライドを処理
                for i, (image, note) in enumerate(zip(images, speaker_notes)):
                    slide_num = i + 1

                    if note:
                        # 発表者ノートがある場合
                        if progress_callback:
                            progress_callback(
                                f"スライド {slide_num}/{total_slides} - 音声生成中...",
                                i + 1, total_steps
                            )

                        # 音声生成
                        audio_path = os.path.join(temp_dir, f"slide_{slide_num:03d}.wav")
                        self.text_to_speech(note, audio_path)

                        # 動画作成
                        video_path = os.path.join(temp_dir, f"slide_{slide_num:03d}.mp4")
                        self.create_slide_video(image, audio_path, video_path)
                    else:
                        # 発表者ノートがない場合は3秒の無音動画
                        if progress_callback:
                            progress_callback(
                                f"スライド {slide_num}/{total_slides} - 無音動画作成中...",
                                i + 1, total_steps
                            )

                        video_path = os.path.join(temp_dir, f"slide_{slide_num:03d}.mp4")
                        self.create_slide_video(image, None, video_path, duration=3.0)

                    video_paths.append(video_path)

                    # メモリ解放
                    del image

                # 4. 動画を結合
                if progress_callback:
                    progress_callback("動画結合中...", total_slides + 1, total_steps)

                self.combine_videos(video_paths, output_path)

            if progress_callback:
                progress_callback("完了", total_steps, total_steps)

            return True, f"講義動画生成完了！ {total_slides}枚のスライドを処理しました。"

        except Exception as e:
            return False, f"エラーが発生しました: {str(e)}"


if __name__ == "__main__":
    import sys

    if len(sys.argv) < 4:
        print("使用方法: python generate_video.py <input.pptx> <slides.zip> <output.mp4>")
        print("ZIPファイルには スライド1.jpeg, スライド2.jpeg... の形式で画像を格納してください")
        sys.exit(1)

    pptx_file = sys.argv[1]
    zip_file = sys.argv[2]
    output_file = sys.argv[3]

    api_key = os.environ.get("GOOGLE_API_KEY")
    if not api_key:
        print("エラー: 環境変数GOOGLE_API_KEYを設定してください")
        sys.exit(1)

    def print_progress(step_name, current, total):
        print(f"[{current}/{total}] {step_name}")

    generator = VideoGenerator(api_key)
    success, message = generator.generate_video(
        pptx_file, zip_file, output_file, progress_callback=print_progress
    )

    print(f"\n結果: {message}")
    if success:
        print(f"出力ファイル: {output_file}")
