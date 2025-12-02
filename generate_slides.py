import json
from pptx import Presentation


def create_slides_from_json(template_path, json_path, output_path):
    """
    JSONファイルから講義スライドを生成
    
    Args:
        template_path: テンプレートPPTXファイルのパス
        json_path: 講義内容を定義したJSONファイルのパス
        output_path: 出力PPTXファイルのパス
    """
    # JSONファイルを読み込む
    with open(json_path, 'r', encoding='utf-8') as f:
        content = json.load(f)
    
    # プレゼンテーションを読み込む
    prs = Presentation(template_path)
    
    # レイアウトを取得
    main_template = prs.slide_masters[0].slide_layouts[0]  # main_template
    start_template = prs.slide_masters[0].slide_layouts[1]  # start
    end_template = prs.slide_masters[0].slide_layouts[2]    # end
    
    # 1枚目: start (単元名)
    slide_start = prs.slides.add_slide(start_template)
    for placeholder in slide_start.placeholders:
        if placeholder.placeholder_format.idx == 13:
            placeholder.text = content['title']
    
    # 2枚目: main_template (アジェンダページ)
    slide_agenda = prs.slides.add_slide(main_template)
    agenda_text = "\n".join([f"{i+1}. {item}" for i, item in enumerate(content['agenda'])])
    for placeholder in slide_agenda.placeholders:
        if placeholder.placeholder_format.idx == 13:
            placeholder.text = "アジェンダ"
        elif placeholder.placeholder_format.idx == 14:
            placeholder.text = agenda_text
    
    # 3枚目以降: 各アジェンダの講義ページ
    for agenda_idx, agenda_title in enumerate(content['agenda']):
        if agenda_title in content['main']:
            # 各アジェンダの複数スライドを生成
            slides_content = content['main'][agenda_title]
            
            for slide_content in slides_content:
                slide = prs.slides.add_slide(main_template)
                for placeholder in slide.placeholders:
                    if placeholder.placeholder_format.idx == 13:
                        placeholder.text = f"{agenda_idx + 1}. {agenda_title}"
                    elif placeholder.placeholder_format.idx == 14:
                        placeholder.text = slide_content
    
    # 最後: end
    slide_end = prs.slides.add_slide(end_template)
    
    # 保存
    prs.save(output_path)
    
    # 統計情報を出力
    total_slides = len(prs.slides)
    print(f"スライド生成完了!")
    print(f"- タイトル: {content['title']}")
    print(f"- アジェンダ数: {len(content['agenda'])}")
    print(f"- 総スライド数: {total_slides}枚")
    print(f"  - start: 1枚")
    print(f"  - アジェンダ: 1枚")
    
    # 各アジェンダのスライド数を集計
    for agenda_idx, agenda_title in enumerate(content['agenda']):
        if agenda_title in content['main']:
            slide_count = len(content['main'][agenda_title])
            print(f"  - {agenda_idx + 1}. {agenda_title}: {slide_count}枚")
    
    print(f"  - end: 1枚")
    print(f"\n出力先: {output_path}")


if __name__ == "__main__":
    # 実行例
    template_path = '/mnt/user-data/uploads/スライト_テンフ_レ.pptx'
    json_path = '/home/claude/lecture_content.json'
    output_path = '/mnt/user-data/outputs/講義スライド_JSON版.pptx'
    
    create_slides_from_json(template_path, json_path, output_path)
